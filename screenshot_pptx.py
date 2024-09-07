from pptx import Presentation
import mahotas as mh

import subprocess
from time import sleep
import tempfile

# Key to advance slides
ADVANCE_SLIDE_KEY = 'J'
# Output file name
DEF_OUTPUT_NAME = 'Snowman.pptx'
# Time to wait between slides (in seconds)
DEF_SLIDE_DELAY = 1

def get_screenshot(fname):
    subprocess.check_call(['scrot', '--monitor', '0', fname])
    return fname


def advance_slide():
    subprocess.check_call(['xdotool', 'key', ADVANCE_SLIDE_KEY])


def parse_args(args):
    import argparse
    parser = argparse.ArgumentParser(formatter_class=argparse.ArgumentDefaultsHelpFormatter,
                                    description='Create PPTX from any presentation (by screenshotting)')
    parser.add_argument('--delay',
                        required=False,
                        type=int,
                        default=DEF_SLIDE_DELAY,
                        action='store',
                        dest='slide_delay')
    parser.add_argument('-o', '--output',
                        required=False,
                        default=DEF_OUTPUT_NAME,
                        dest='output_name')
    return parser.parse_args(args)


def main(args=None):
    if args is None:
        from sys import argv
        args = argv[1:]
    opts = parse_args(args)

    print('Creating slides')
    print(f'Will simulate pressing {ADVANCE_SLIDE_KEY} to advance slides')
    print('Press Ctrl+C to stop')
    for i in range(5):
        print(f'Starting in {5-i} seconds')
        sleep(1)

    prs = Presentation('empty_16x10.pptx')
    blank_slide_layout = prs.slide_layouts[0]

    prev = None
    while True:
        with tempfile.TemporaryDirectory() as tdir:
            cur = get_screenshot(tdir + 'test.png')
            im = mh.imread(cur)
            if prev is not None:
                if (prev == im).all():
                    print('End of presentation detected (no change in screenshot)')
                    break
                slide = prs.slides.add_slide(blank_slide_layout)
            else:
                slide = prs.slides[0]
            slide.shapes.add_picture(cur, 0, 0, height=prs.slide_height)
            prev = im

        advance_slide()
        sleep(opts.slide_delay)

    prs.save(opts.output_name)
    print(f'Saved to {opts.output_name}')

if __name__ == '__main__':
    main()
